#!/usr/bin/env python3
"""Convert PDF / Office documents to Markdown with streaming, parallel processing.

Features:
  - Streaming batch processing for 300+ page PDFs (memory-efficient)
  - Parallel OCR and image rendering (ThreadPoolExecutor)
  - Tesseract OCR fallback when pymupdf built-in OCR is unavailable
  - Office document support (DOCX / XLSX / PPTX) via native Python libs
  - macOS textutil fallback for DOC / ODT
  - Per-page error recovery (continues on failure, reports at end)
  - Progress reporting to stderr (JSON lines, parsable by PHP)
"""

from __future__ import annotations

import argparse
import json
import math
import os
import re
import subprocess
import sys
import time
from collections import Counter
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path
from typing import Any, Dict, List, Optional, Set, Tuple

# ── Optional third-party dependencies ────────────────────────────────────────

try:
    import pymupdf
    import pymupdf4llm

    HAS_PYMUPDF = True
except ImportError:
    pymupdf = pymupdf4llm = None
    HAS_PYMUPDF = False

try:
    import docx as docx_module

    HAS_DOCX = True
except ImportError:
    docx_module = None
    HAS_DOCX = False

try:
    import openpyxl

    HAS_OPENPYXL = True
except ImportError:
    openpyxl = None
    HAS_OPENPYXL = False

try:
    import pptx as pptx_module

    HAS_PPTX = True
except ImportError:
    pptx_module = None
    HAS_PPTX = False

try:
    from PIL import Image

    HAS_PIL = True
except ImportError:
    Image = None
    HAS_PIL = False

try:
    import pygments
    from pygments.lexers import guess_lexer, get_lexer_by_name
    from pygments.util import ClassNotFound as PygmentsClassNotFound

    HAS_PYGMENTS = True
except ImportError:
    pygments = None
    guess_lexer = None
    get_lexer_by_name = None
    PygmentsClassNotFound = Exception
    HAS_PYGMENTS = False

# ── Constants ────────────────────────────────────────────────────────────────

DEFAULT_BATCH_SIZE = 50
OFFICE_EXTENSIONS = {'.docx', '.doc', '.xlsx', '.pptx', '.odt', '.ods', '.odp'}
TEXTUTIL_FORMATS = {'doc', 'docx', 'odt', 'rtf'}
SUPPORTED_EXTENSIONS = {'.pdf'} | OFFICE_EXTENSIONS

# ── Exceptions ───────────────────────────────────────────────────────────────


class ConversionError(Exception):
    pass


# ── CLI ──────────────────────────────────────────────────────────────────────


def parse_args(argv: Optional[List[str]] = None) -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument('--input', required=True, type=Path)
    parser.add_argument('--output-dir', required=True, type=Path)
    parser.add_argument('--public-prefix', required=True)
    parser.add_argument('--max-pages', type=int, default=0,
                        help='Max pages (0 = unlimited, default)')
    parser.add_argument('--image-dpi', type=int, default=144)
    parser.add_argument('--visual-dpi', type=int, default=110)
    parser.add_argument('--ocr-language', default='fra+eng')
    parser.add_argument('--batch-size', type=int, default=DEFAULT_BATCH_SIZE)
    parser.add_argument('--parallel', type=int, default=0,
                        help='Workers (0 = auto = CPU count)')
    parser.add_argument('--skip-parallel', action='store_true',
                        help='Disable parallel processing')
    return parser.parse_args(argv)

# ── Progress helper ──────────────────────────────────────────────────────────


def log_progress(stage: str, current: int, total: int, message: str = '') -> None:
    data = {'stage': stage, 'current': current, 'total': total, 'message': message}
    print(json.dumps(data, ensure_ascii=False), file=sys.stderr, flush=True)

# ── Path helpers ─────────────────────────────────────────────────────────────


def rewrite_image_paths(markdown: str, output_dir: Path, prefix: str) -> str:
    normalized = output_dir.resolve().as_posix().rstrip('/')
    markdown = markdown.replace(normalized + '/', prefix.rstrip('/') + '/')
    markdown = markdown.replace(str(output_dir).rstrip('/') + '/', prefix.rstrip('/') + '/')
    return markdown

# ── File type detection ──────────────────────────────────────────────────────


def detect_file_type(path: Path) -> str:
    ext = path.suffix.lower()
    if ext == '.pdf':
        return 'pdf'
    if ext in OFFICE_EXTENSIONS:
        return 'office'
    return 'unknown'


def office_subtype(path: Path) -> str:
    return {'.docx': 'docx', '.doc': 'doc', '.xlsx': 'xlsx',
            '.xls': 'xls', '.pptx': 'pptx', '.ppt': 'ppt',
            '.odt': 'odt', '.ods': 'ods', '.odp': 'odp'}.get(path.suffix.lower(), 'unknown')

# ── Utility functions ────────────────────────────────────────────────────────


def render_page(page: pymupdf.Page, path: Path, dpi: int) -> None:
    scale = dpi / 72
    pixmap = page.get_pixmap(matrix=pymupdf.Matrix(scale, scale), alpha=False)
    pixmap.save(str(path))


def normalize_boundary_line(line: str) -> str:
    line = re.sub(r'^#{1,6}\s+', '', line.strip())
    line = re.sub(r'[*_`>#|]', '', line)
    line = re.sub(r'\d+', '#', line)
    line = re.sub(r'\s+', ' ', line)
    return line.strip(' -–—·').casefold()


def is_page_number(line: str) -> bool:
    normalized = re.sub(r'[\u200b-\u200f\ufeff]', '', line).strip()
    normalized = re.sub(r'[*_`#>]', '', normalized).strip()
    return bool(re.fullmatch(
        r'(?:page\s*)?[-–—]?\s*\d+\s*(?:[/|]\s*\d+|(?:sur|of)\s+\d+)?\s*[-–—]?',
        normalized, flags=re.IGNORECASE,
    ))


def page_image_coverage(page: pymupdf.Page) -> float:
    page_area = max(page.rect.width * page.rect.height, 1)
    image_area = 0.0

    for image in page.get_image_info():
        bbox = image.get('bbox')
        if bbox:
            rect = pymupdf.Rect(bbox) & page.rect
            image_area += max(rect.width, 0) * max(rect.height, 0)

    return min(image_area / page_area, 1.0)


def page_requires_visual_reference(page: pymupdf.Page, text: str) -> bool:
    """Keep a visual reference only when the page cannot be represented reliably as text."""
    if len(text.strip()) < 40:
        return True
    return page_image_coverage(page) >= 0.08 or len(page.get_drawings()) >= 8


def classify_document_profile(page_profiles: List[Dict[str, int]]) -> Dict[str, Any]:
    """Select an extraction strategy from cheap, dependency-free page metrics."""
    total = max(len(page_profiles), 1)
    scanned_pages = sum(1 for profile in page_profiles if profile['is_scanned'])
    complex_pages = sum(1 for profile in page_profiles if profile['is_complex'])
    image_pages = sum(1 for profile in page_profiles if profile['image_count'] > 0)
    scanned_ratio = scanned_pages / total
    complex_ratio = complex_pages / total

    if scanned_ratio >= 0.8:
        document_type = 'scanned'
        strategy = 'ocr-first'
    elif scanned_pages > 0:
        document_type = 'hybrid'
        strategy = 'mixed-text-ocr'
    elif complex_ratio >= 0.35:
        document_type = 'complex'
        strategy = 'layout-aware'
    else:
        document_type = 'native'
        strategy = 'text-first'

    return {
        'document_type': document_type,
        'strategy': strategy,
        'scanned_pages': scanned_pages,
        'complex_pages': complex_pages,
        'image_pages': image_pages,
    }


def inspect_pdf(doc: pymupdf.Document) -> Tuple[List[Dict[str, int]], Dict[str, Any]]:
    profiles: List[Dict[str, int]] = []

    for page in doc:
        text_length = len(page.get_text('text').strip())
        image_count = len(page.get_images(full=True))
        drawing_count = len(page.get_drawings())
        image_coverage = page_image_coverage(page)
        is_scanned = int(text_length < 40 and image_coverage >= 0.5)
        is_complex = int(image_coverage >= 0.08 or drawing_count >= 8)
        profiles.append({
            'text_length': text_length,
            'image_count': image_count,
            'drawing_count': drawing_count,
            'image_coverage_percent': round(image_coverage * 100),
            'is_scanned': is_scanned,
            'is_complex': is_complex,
        })

    return profiles, classify_document_profile(profiles)


def clean_markdown_page(markdown: str, repeated_boundaries: Set[str]) -> str:
    cleaned_lines: List[str] = []
    for line in markdown.splitlines():
        if is_page_number(line):
            continue
        if normalize_boundary_line(line) in repeated_boundaries:
            continue
        cleaned_lines.append(line.rstrip())
    cleaned = '\n'.join(cleaned_lines)
    cleaned = re.sub(r'!\[\s*\]\([^\n)]+\)', '', cleaned)
    cleaned = re.sub(r'[ \t]+\n', '\n', cleaned)
    cleaned = re.sub(r'\n{3,}', '\n\n', cleaned)
    return cleaned.strip()


def detect_repeated_boundaries(markdown_pages: List[str]) -> Set[str]:
    """Streaming boundary detection — samples pages instead of full O(n²)."""
    if len(markdown_pages) < 2:
        return set()

    occurrences: Counter = Counter()
    sample_count = min(len(markdown_pages), 120)

    if len(markdown_pages) > 120:
        indices: Set[int] = set()
        indices.update(range(0, min(40, len(markdown_pages))))
        indices.update(range(max(0, len(markdown_pages) - 30), len(markdown_pages)))
        step = max(1, (len(markdown_pages) - 70) // 50)
        indices.update(range(40, len(markdown_pages) - 30, step))
        sample_indices = sorted(indices)[:sample_count]
    else:
        sample_indices = range(len(markdown_pages))

    for idx in sample_indices:
        lines = [ln for ln in markdown_pages[idx].splitlines() if ln.strip()]
        candidates = {
            normalize_boundary_line(ln)
            for ln in [*lines[:3], *lines[-3:]]
            if normalize_boundary_line(ln)
        }
        occurrences.update(candidates)

    threshold = max(2, math.ceil(sample_count * 0.4))
    return {ln for ln, cnt in occurrences.items() if cnt >= threshold}


def repeated_boundary_lines(markdown_pages: List[str]) -> Set[str]:
    """Backward-compatible public name used by cleanup tests and callers."""
    return detect_repeated_boundaries(markdown_pages)

# ── OCR ──────────────────────────────────────────────────────────────────────


def ocr_via_pymupdf(page: pymupdf.Page, language: str) -> str:
    textpage = page.get_textpage_ocr(language=language, dpi=200, full=True)
    return page.get_text('text', textpage=textpage).strip()


def ocr_via_tesseract(page: pymupdf.Page, language: str) -> str:
    scale = 200 / 72
    pix = page.get_pixmap(matrix=pymupdf.Matrix(scale, scale), alpha=False)
    img_bytes = pix.tobytes('png')

    language_candidates = list(dict.fromkeys([language, *language.split('+'), 'eng']))
    for candidate in language_candidates:
        try:
            result = subprocess.run(
                ['tesseract', 'stdin', 'stdout', '-l', candidate, '--psm', '3'],
                input=img_bytes,
                capture_output=True,
                timeout=120,
            )
            if result.returncode == 0:
                return result.stdout.decode('utf-8').strip()
        except (FileNotFoundError, subprocess.TimeoutExpired):
            break

    return ''


def try_ocr(page: pymupdf.Page, language: str, page_number: int) -> Tuple[str, bool, str]:
    """Try pymupdf OCR first, then Tesseract fallback.
    Returns (ocr_text, success, method_used).
    """
    method = ''
    try:
        text = ocr_via_pymupdf(page, language)
        if text:
            method = 'pymupdf'
            return text, True, method
    except Exception:
        pass

    try:
        text = ocr_via_tesseract(page, language)
        if text:
            method = 'tesseract'
            return text, True, method
    except Exception:
        pass

    return '', False, ''

# ── Image extraction from Office documents ────────────────────────────────────


def extract_image_blob(output_dir: Path, blob: bytes, ext: str, counter: int, prefix: str) -> str:
    """Save an image blob to the output directory and return the markdown reference."""
    filename = f'office-img-{counter:04d}{ext}'
    (output_dir / filename).write_bytes(blob)
    return f'![Image]({prefix.rstrip("/")}/{filename})'


def extract_docx_images(doc: Any, output_dir: Path, prefix: str) -> Tuple[List[str], int]:
    """Extract embedded images from a python-docx Document."""
    refs: List[str] = []
    count = 0
    seen: Set[str] = set()
    for rel in doc.part.rels.values():
        if 'image' in (rel.reltype or ''):
            try:
                image = rel.target_part
                blob = image.blob
                key = str(hash(blob))
                if key in seen:
                    continue
                seen.add(key)
                ext = Path(image.partname).suffix or '.png'
                ref = extract_image_blob(output_dir, blob, ext, count, prefix)
                refs.append(ref)
                count += 1
            except Exception:
                continue
    return refs, count


def extract_pptx_images(slide_shape: Any, output_dir: Path, prefix: str, counter: int) -> Tuple[List[str], int]:
    """Extract images from a python-pptx slide shape."""
    refs: List[str] = []
    count = 0
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        if slide_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            blob = slide_shape.image.blob
            ext = '.' + (slide_shape.image.content_type.split('/')[-1] or 'png')
            ref = extract_image_blob(output_dir, blob, ext, counter, prefix)
            refs.append(ref)
            count += 1
    except Exception:
        pass
    return refs, count

# ── Code block detection ─────────────────────────────────────────────────────

# Quick keyword sets for first-pass heuristics (before calling Pygments)
CODE_HINTS: Dict[str, List[str]] = {
    'python': ['def ', 'class ', 'import ', 'from ', 'elif ', 'else:', 'try:', 'except ', 'finally:', 'yield ', 'lambda ', 'self.', '__init__', '__str__', 'if __name__', 'return ', 'range(', 'len(', 'print(', 'True', 'False', 'None', 'in ', 'not ', 'and ', 'or ', 'pass', 'raise ', 'with ', 'as ', 'async def', 'await '],
    'php': ['<?php', 'function ', 'public ', 'private ', 'protected ', 'static ', '->', 'namespace ', 'use ', 'echo ', 'array_map', 'array_filter', '$this', 'new ', 'return ', 'true', 'false', 'null'],
    'javascript': ['function ', 'const ', 'let ', 'var ', '=> ', 'console.', 'document.', 'window.', 'async ', 'await ', 'null', 'undefined', 'prototype.', 'new ', 'return ', 'this.', 'true', 'false', 'class ', 'extends '],
    'java': ['public class', 'private ', 'protected ', 'static ', 'void ', 'String ', 'int ', 'new ', 'extends ', 'implements ', '@Override', 'System.out', 'return ', 'true', 'false', 'null', 'import java'],
    'cpp': ['#include', 'int main', 'std::', 'cout', 'cin', 'namespace std', 'public:', 'private:', 'virtual ', 'template', '->', 'const ', 'auto ', 'true', 'false', 'nullptr'],
    'csharp': ['using ', 'namespace ', 'class ', 'public ', 'private ', 'void ', 'string ', 'int ', 'var ', 'async ', 'await ', 'Console.', 'return ', 'true', 'false', 'null'],
    'html': ['<!DOCTYPE', '<html', '<head', '<body', '<div', '<span', '<table', '<form', '<input', '<!DOCTYPE html', '<a ', '<img', '<script', '<style'],
    'css': ['{', 'px;', '!important', '@media', 'display:', 'margin:', 'padding:', 'flex:', 'grid:', ':hover', ':before', ':after', '#id', '.class'],
    'sql': ['SELECT ', 'FROM ', 'WHERE ', 'INSERT ', 'UPDATE ', 'DELETE ', 'CREATE TABLE', 'ALTER ', 'JOIN ', 'GROUP BY', 'ORDER BY', 'HAVING', 'VALUES ', 'SET ', 'INNER JOIN', 'LEFT JOIN'],
    'go': ['func ', 'package ', 'import (', 'defer ', 'go ', 'chan ', 'interface{', 'struct {', ':=', 'nil', 'true', 'false', 'make(', 'append(', 'len('],
    'rust': ['fn ', 'let ', 'mut ', 'impl ', 'trait ', 'pub ', 'struct ', 'enum ', 'match ', 'println!', 'unwrap()', 'Some(', 'None', 'Ok(', 'Err(', 'String', 'vec!'],
    'ruby': ['def ', 'class ', 'end', 'do |', '->', 'puts ', 'require ', 'attr_', 'private', 'protected', 'true', 'false', 'nil', 'each ', 'map ', 'select '],
    'swift': ['func ', 'var ', 'let ', 'class ', 'struct ', 'enum ', 'guard ', 'defer ', '-> ', '?? ', 'true', 'false', 'nil', 'import '],
    'kotlin': ['fun ', 'val ', 'var ', 'class ', 'object ', 'override ', 'private ', 'import ', 'println', 'true', 'false', 'null', 'when ', 'sealed '],
}


def looks_like_code(text: str) -> bool:
    """Heuristic check: is this block likely source code rather than prose?"""
    lines = [ln for ln in text.split('\n') if ln.strip()]
    if len(lines) < 2:
        return False

    # Code typically has indented lines
    indented = sum(1 for ln in lines if ln.startswith(('    ', '\t', '  ')))
    has_indentation = indented >= 2 or (indented > 0 and indented / len(lines) > 0.3)

    # Code often has special characters on multiple lines
    has_braces = sum(1 for ln in lines if '{' in ln or '}' in ln or '(' in ln) >= 2
    has_semicolons = sum(1 for ln in lines if ';' in ln) >= 2
    has_equals = sum(1 for ln in lines if '=' in ln and not ln.strip().startswith('=')) >= 2

    # Short lines (code tends to have shorter average lines than prose)
    avg_len = sum(len(ln) for ln in lines) / len(lines)

    # Prose usually has longer lines with spaces, code has more symbols
    symbol_ratio = sum(1 for c in text if c in '{}();=<>!&|') / max(len(text), 1)

    # Check for programming keywords
    kw_hits = sum(1 for kws in CODE_HINTS.values() for kw in kws if kw in text)

    score = 0
    if has_indentation:
        score += 2
    if has_braces:
        score += 2
    if has_semicolons:
        score += 1
    if has_equals:
        score += 1
    if avg_len < 80:
        score += 1
    if symbol_ratio > 0.03:
        score += 1
    if kw_hits >= 3:
        score += 2
    elif kw_hits >= 1:
        score += 1

    return score >= 4


def detect_code_language(code: str) -> str:
    """Detect programming language — heuristic first, Pygments as fallback."""
    best_lang = ''
    best_score = 0
    for lang, keywords in CODE_HINTS.items():
        score = sum(1 for kw in keywords if kw in code)
        if score > best_score:
            best_score = score
            best_lang = lang

    # If heuristics give a clear winner (3+ keyword hits), trust it
    if best_score >= 3:
        return best_lang

    if HAS_PYGMENTS:
        try:
            lexer = guess_lexer(code)
            if lexer.aliases:
                pyg_lang = lexer.aliases[0]
                # Validate Pygments result against our hints
                if best_score == 0:
                    return pyg_lang
                # If both agree, use it; if not, trust heuristics
                if pyg_lang == best_lang or pyg_lang in {'python', 'javascript', 'php', 'java', 'cpp', 'csharp'}:
                    return pyg_lang
        except (PygmentsClassNotFound, Exception):
            pass

    return best_lang if best_score >= 1 else ''


def annotate_code_blocks(markdown: str) -> str:
    """Find code blocks in markdown and wrap them with ```language annotations.

    Uses a two-pass approach:
    1. Identify high-density code regions (including isolated comment lines)
    2. Wrap those regions in fenced code blocks
    """
    lines = markdown.split('\n')
    n = len(lines)

    # Pass 1: classify each line
    # -1 = fence, 0 = non-code, 1 = code, 2 = blank
    classified: List[int] = []
    for line in lines:
        s = line.strip()
        if s.startswith('```') or s.startswith('~~~'):
            classified.append(-1)
        elif not s:
            classified.append(2)
        elif is_code_line_candidate(line):
            classified.append(1)
        else:
            classified.append(0)

    # Pass 2: find code regions (contiguous 1s with max 1 blank, include isolated 0s between 1s)
    regions: List[Tuple[int, int]] = []
    i = 0
    while i < n:
        if classified[i] == 1:
            start = i
            # Expand until we hit 2+ blanks or fence or non-code non-adjacent
            j = i + 1
            blanks = 0
            code_adjacent = True
            while j < n and classified[j] != -1:
                if classified[j] == 1:
                    blanks = 0
                    code_adjacent = True
                    j += 1
                elif classified[j] == 2:
                    blanks += 1
                    if blanks > 1:
                        break
                    j += 1
                elif classified[j] == 0:
                    # If we've seen code before in this region, include isolated non-code
                    if j < n - 1 and classified[j + 1] == 1 and blanks == 0:
                        # Non-code line sandwiched between code → include it
                        blanks = 0
                        j += 1
                        continue
                    # Otherwise stop
                    break
            regions.append((start, j))
            i = j
        else:
            i += 1

    # Merge overlapping or adjacent regions
    if regions:
        merged = [regions[0]]
        for r in regions[1:]:
            prev = merged[-1]
            if r[0] - prev[1] <= 1:
                merged[-1] = (prev[0], r[1])
            else:
                merged.append(r)
        regions = merged

    # Pass 3: build output, wrapping code regions
    result: List[str] = []
    region_idx = 0
    i = 0
    while i < n:
        if region_idx < len(regions) and i == regions[region_idx][0]:
            start, end = regions[region_idx]
            code_text = '\n'.join(line for line in lines[start:end] if line.strip()).strip()
            if looks_like_code(code_text):
                lang = detect_code_language(code_text)
                ann = f'```{lang}\n{code_text}\n```' if lang else f'```\n{code_text}\n```'
                result.append(ann)
            else:
                result.append('\n'.join(lines[start:end]))
            i = end
            region_idx += 1
        else:
            result.append(lines[i])
            i += 1

    return '\n'.join(result)


def is_code_line_candidate(line: str) -> bool:
    """Quick check if a line looks like it could be part of a code block."""
    # Check indent BEFORE strip
    has_indent = line.startswith(('    ', '\t', '  ')) if line else False

    s = line.strip()
    if not s:
        return False

    # Only exclude heading markers if NOT indented (indented '#' is a code comment)
    if not has_indent:
        # Heading: short text after #, typically capitalised, no code chars
        heading_match = re.match(r'^#{1,6}\s+(.+)', s)
        if heading_match:
            heading_text = heading_match.group(1)
            # A heading is typically < 60 chars, no special chars, starts with capital
            is_heading = (
                len(heading_text) < 60
                and not any(c in heading_text for c in '{}();=<>[]')
                and (heading_text[0].isupper() if heading_text else True)
            )
            if is_heading:
                return False
        if re.match(r'^[\*\-\+]\s', s) or re.match(r'^\d+\.\s', s):
            return False
        if s.startswith('!['):
            return False
        if '|' in s and s.count('|') >= 2:
            return False

    if s.startswith('```') or s.startswith('~~~'):
        return False

    has_code_chars = any(c in s for c in '{}();=<>&|#')
    has_keywords = any(any(kw in s for kw in kws) for kws in CODE_HINTS.values())
    short_with_chars = has_code_chars and len(s) < 60
    return has_keywords or has_indent or short_with_chars


# ── Equation / math detection ────────────────────────────────────────────────

MATH_UNICODE = set('∫∑∏√π∞∂∆θλμσΩωαβγδεζηικνοπρςτυφχψ∇∈∉∋∌⊂⊃⊄⊅∪∩∅∧∨¬⇒⇔∀∃∴∵∘∙⊕⊗⊥⊢⊣⊤⊥')
MATH_PATTERNS = re.compile(
    r'(?:'
    r'\\frac\{[^}]+\}\{[^}]+\}|'
    r'\\sqrt(?:\[[^\]]+\])?\{[^}]+\}|'
    r'\\sum(?:_{[^}]+})?(?:\^{[^}]+})?|'
    r'\\int(?:_{[^}]+})?(?:\^{[^}]+})?|'
    r'\\lim|\\log|\\ln|\\sin|\\cos|\\tan|\\partial|'
    r'\\rightarrow|\\leftarrow|\\Rightarrow|\\Leftarrow|'
    r'\\alpha|\\beta|\\gamma|\\delta|\\epsilon|\\theta|\\lambda|\\pi|\\sigma|\\omega|'
    r'\\pm|\\mp|\\times|\\div|\\leq|\\geq|\\approx|\\neq|'
    r'\{[^}]*\}^{[^}]*|'                          # a^{b}
    r'\{[^}]*\}_{[^}]*|'                          # a_{b}
    r'(?<![a-zA-Z])[a-zA-Z]\s*=\s*[a-zA-Z]|'      # a = b
    r'(?<!\w)\d+\.\d+(?!\w)|'                     # decimal numbers
    r'(?<!\w)[a-zA-Z]\^[0-9](?!\w)|'              # x^2
    r'[a-zA-Z]\_[a-zA-Z]|'                         # a_b
    r'[-+]?\d+\s*[\+\-\*/]\s*[-+]?\d+'            # 2 + 2 arithmetic
    r')'
)


def contains_math(text: str) -> bool:
    """Check if a text block contains mathematical content."""
    if MATH_PATTERNS.search(text):
        return True
    if any(ch in text for ch in MATH_UNICODE):
        return True
    return False


def wrap_equations(markdown: str) -> str:
    """Wrap detected math expressions in $...$ or $$...$$, skipping code blocks."""
    lines = markdown.split('\n')
    result: List[str] = []
    inside_code = False
    i = 0
    while i < len(lines):
        line = lines[i]
        stripped = line.strip()

        # Track code fences
        if stripped.startswith('```') or stripped.startswith('~~~'):
            inside_code = not inside_code
            result.append(line)
            i += 1
            continue

        if inside_code:
            result.append(line)
            i += 1
            continue

        # Skip if already in math mode
        if stripped.startswith('$$') or stripped.startswith('$') or stripped.startswith(r'\['):
            result.append(line)
            i += 1
            continue

        # Skip headings, list items, images, tables
        if re.match(r'^#{1,6}\s', stripped) or re.match(r'^[\*\-\+]\s', stripped):
            result.append(line)
            i += 1
            continue
        if re.match(r'^\d+\.\s', stripped) or stripped.startswith('!['):
            result.append(line)
            i += 1
            continue
        if '|' in stripped and stripped.count('|') >= 2:
            result.append(line)
            i += 1
            continue

        # Only apply equation detection on standalone prose lines that look like math
        if contains_math(stripped) and len(stripped) > 5 and not looks_like_code(stripped):
            math_lines = [stripped]
            j = i + 1
            while (j < len(lines) and contains_math(lines[j].strip())
                   and len(lines[j].strip()) > 5
                   and not lines[j].strip().startswith('```')
                   and not looks_like_code(lines[j].strip())):
                math_lines.append(lines[j].strip())
                j += 1

            if len(math_lines) >= 2:
                result.append('$$')
                result.extend(math_lines)
                result.append('$$')
                i = j
                continue
            else:
                # Only wrap as equation if it's clearly math (not prose with =)
                if MATH_UNICODE & set(stripped) or re.search(r'\\[a-z]+', stripped):
                    result.append(f'${stripped}$')
                    i += 1
                    continue

        result.append(line)
        i += 1

    return '\n'.join(result)


def post_process_markdown(markdown: str) -> str:
    """Apply code detection and equation wrapping to assembled markdown."""
    markdown = annotate_code_blocks(markdown)
    markdown = wrap_equations(markdown)
    return markdown

# ── Office extraction ────────────────────────────────────────────────────────


def extract_docx(path: Path, output_dir: Path, prefix: str) -> Tuple[str, List[str], int]:
    if not HAS_DOCX:
        raise ConversionError('python-docx non disponible (pip install python-docx)')
    doc = docx_module.Document(str(path))
    parts: List[str] = []
    warnings: List[str] = []
    image_refs: List[str] = []

    # Extract images first
    img_refs, img_count = extract_docx_images(doc, output_dir, prefix)

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style = (para.style.name or '').lower()
        level = 0
        for i in range(1, 7):
            if str(i) in style or f'h{i}' in style or f'titre {i}' in style:
                level = i
                break
        if level:
            parts.append(f'{"#" * level} {text}')
        elif any(r.bold for r in para.runs if r.bold):
            bold = ' '.join(r.text for r in para.runs if r.bold)
            parts.append(f'**{bold}**')
        elif any(r.italic for r in para.runs if r.italic):
            italic = ' '.join(r.text for r in para.runs if r.italic)
            parts.append(f'*{italic}*')
        else:
            parts.append(text)

    # Insert image references where they appear in the document flow
    # python-docx doesn't expose image position in text flow, so append at end
    if img_refs:
        parts.append('')
        parts.extend(img_refs)

    for table in doc.tables:
        parts.append('')
        for row in table.rows:
            cells = [cell.text.strip().replace('\n', ' ') for cell in row.cells]
            parts.append('| ' + ' | '.join(cells) + ' |')
        parts.append('')
    return '\n\n'.join(parts).strip(), warnings, img_count


def extract_xlsx(path: Path) -> Tuple[str, List[str]]:
    if not HAS_OPENPYXL:
        raise ConversionError('openpyxl non disponible (pip install openpyxl)')
    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    parts: List[str] = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        parts.append(f'## {sheet_name}\n')
        rows: List[List[str]] = []
        for row in ws.iter_rows(values_only=True):
            rows.append([str(c) if c is not None else '' for c in row])
        if rows:
            parts.append('| ' + ' | '.join(rows[0]) + ' |')
            parts.append('| ' + ' | '.join(['---'] * len(rows[0])) + ' |')
            for r in rows[1:]:
                parts.append('| ' + ' | '.join(r) + ' |')
    wb.close()
    return '\n'.join(parts).strip(), []


def extract_pptx(path: Path, output_dir: Path, prefix: str) -> Tuple[str, List[str], int, int]:
    if not HAS_PPTX:
        raise ConversionError('python-pptx non disponible (pip install python-pptx)')
    prs = pptx_module.Presentation(str(path))
    parts: List[str] = []
    slide_count = len(list(prs.slides))
    total_images = 0
    for num, slide in enumerate(prs.slides, 1):
        slide_parts: List[str] = [f'## Diapositive {num}']
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        slide_parts.append(t)
            # Extract images from shapes
            refs, cnt = extract_pptx_images(shape, output_dir, prefix, total_images)
            total_images += cnt
            slide_parts.extend(refs)
        parts.append('\n'.join(slide_parts))
    return '\n\n'.join(parts).strip(), [], slide_count, total_images


def extract_via_textutil(path: Path) -> Tuple[str, List[str]]:
    ext = path.suffix.lower().lstrip('.')
    if ext not in TEXTUTIL_FORMATS:
        raise ConversionError(f'textutil ne supporte pas .{ext}')
    try:
        result = subprocess.run(
            ['textutil', '-convert', 'txt', '-stdout', str(path)],
            capture_output=True, text=True, timeout=60,
        )
        if result.returncode != 0:
            raise ConversionError(f'textutil a échoué: {result.stderr.strip()}')
        return result.stdout.strip(), []
    except FileNotFoundError:
        raise ConversionError('textutil introuvable (macOS uniquement)')
    except subprocess.TimeoutExpired:
        raise ConversionError('textutil a expiré')


def extract_office(args: argparse.Namespace) -> Dict[str, Any]:
    path = args.input.resolve()
    output_dir = args.output_dir.resolve()
    subtype = office_subtype(path)
    markdown = ''
    warnings: List[str] = []
    page_count = 1
    image_count = 0

    if subtype == 'docx':
        markdown, warnings, image_count = extract_docx(path, output_dir, args.public_prefix)
    elif subtype in ('xlsx', 'xls'):
        markdown, warnings = extract_xlsx(path)
        page_count = max(1, markdown.count('## '))
    elif subtype in ('pptx', 'ppt'):
        markdown, warnings, slide_count, image_count = extract_pptx(path, output_dir, args.public_prefix)
        page_count = slide_count
    elif subtype in ('doc', 'odt'):
        markdown, warnings = extract_via_textutil(path)
    else:
        markdown, warnings = extract_via_textutil(path)

    markdown = post_process_markdown(markdown)
    word_count = len(re.findall(r'\b[\wÀ-ÿ\'-]+\b', markdown, flags=re.UNICODE))

    office_warnings = warnings
    if image_count == 0:
        office_warnings = warnings + [
            "Aucune image extraite (les images peuvent ne pas être conservées pour certains formats Office).",
        ]

    return {
        'markdown': markdown,
        'page_count': page_count,
        'word_count': word_count,
        'image_count': image_count,
        'cover_file': '',
        'visual_pages': [],
        'ocr_pages': [],
        'ocr_required_pages': [],
        'warnings': office_warnings,
    }

# ── PDF extraction in batches ────────────────────────────────────────────────


def ocr_pdf_page(input_path: Path, page_number: int, language: str) -> Tuple[int, str, bool, str]:
    """OCR one page using an isolated document handle, safe for parallel workers."""
    worker_doc = pymupdf.open(str(input_path))
    try:
        text, success, method = try_ocr(worker_doc[page_number - 1], language, page_number)
        return page_number, text, success, method
    finally:
        worker_doc.close()


def process_pdf(args: argparse.Namespace) -> Dict[str, Any]:
    input_path = args.input.resolve()
    output_dir = args.output_dir.resolve()
    output_dir.mkdir(parents=True, exist_ok=True)

    if not HAS_PYMUPDF:
        raise ConversionError(
            'PyMuPDF indisponible. Exécutez: '
            'pip install pymupdf pymupdf4llm'
        )

    doc = pymupdf.open(str(input_path))

    if doc.needs_pass:
        raise ConversionError('Le PDF est protégé par un mot de passe.')
    if doc.page_count == 0:
        raise ConversionError('Le PDF ne contient aucune page.')
    if args.max_pages > 0 and doc.page_count > args.max_pages:
        raise ConversionError(
            f'Le PDF contient {doc.page_count} pages; limite = {args.max_pages}.'
        )

    total = doc.page_count
    page_profiles, document_profile = inspect_pdf(doc)
    parallel = min(os.cpu_count() or 4, 4) if args.parallel == 0 else max(args.parallel, 1)
    if args.skip_parallel:
        parallel = 1

    batch_size = min(max(args.batch_size, 1), max(10, total // parallel + 1))
    batch_size = min(batch_size, total)

    log_progress(
        'start',
        0,
        total,
        f'{total} pages, type {document_profile["document_type"]}, '
        f'stratégie {document_profile["strategy"]}, lots de {batch_size}, {parallel} workers',
    )

    # Cover
    cover = 'cover.png'
    render_page(doc[0], output_dir / cover, args.visual_dpi)

    raw_pages: List[str] = []
    scanned: List[int] = []
    ocr_done: List[int] = []
    visual_queue: List[int] = []
    warnings: List[str] = []

    # ── Phase 1: adaptive extraction ─────────────────────────────────────
    if document_profile['document_type'] == 'scanned':
        page_numbers = list(range(1, total + 1))
        ocr_results: Dict[int, Tuple[str, bool, str]] = {}
        log_progress('ocr', 0, total, 'OCR parallèle du document scanné')

        if parallel > 1 and total > 1:
            with ThreadPoolExecutor(max_workers=parallel) as executor:
                futures = {
                    executor.submit(ocr_pdf_page, input_path, page_number, args.ocr_language): page_number
                    for page_number in page_numbers
                }
                for future in as_completed(futures):
                    page_number = futures[future]
                    try:
                        _, text, success, method = future.result()
                        ocr_results[page_number] = (text, success, method)
                    except Exception as error:
                        ocr_results[page_number] = ('', False, '')
                        warnings.append(f'Erreur OCR page {page_number}: {error}')
        else:
            for page_number in page_numbers:
                _, text, success, method = ocr_pdf_page(
                    input_path,
                    page_number,
                    args.ocr_language,
                )
                ocr_results[page_number] = (text, success, method)

        for page_number in page_numbers:
            text, success, method = ocr_results[page_number]
            raw_pages.append(text)
            visual_queue.append(page_number)
            if success:
                ocr_done.append(page_number)
                log_progress('ocr', page_number, total, f'OCR page {page_number} via {method}')
            else:
                scanned.append(page_number)
                log_progress('ocr', page_number, total, f'OCR impossible page {page_number}')
    else:
        extract_images = document_profile['image_pages'] > 0

        for batch_start in range(0, total, batch_size):
            batch_end = min(batch_start + batch_size, total)
            log_progress('extract', batch_end, total, f'Lot {batch_start + 1}-{batch_end}')

            batch_doc = pymupdf.open()
            batch_doc.insert_pdf(doc, from_page=batch_start, to_page=batch_end - 1)

            try:
                chunks = pymupdf4llm.to_markdown(
                    batch_doc,
                    page_chunks=True,
                    write_images=extract_images,
                    image_path=str(output_dir),
                    image_format='png',
                    dpi=args.image_dpi,
                    force_text=True,
                    show_progress=False,
                    margins=(0, 36, 0, 36),
                )
            finally:
                batch_doc.close()

            if not chunks:
                warnings.append(f'Lot {batch_start + 1}-{batch_end}: aucune donnée extraite')
                continue

            for batch_index, chunk in enumerate(chunks):
                page_number = batch_start + batch_index + 1
                if page_number > total:
                    break

                page = doc[page_number - 1]
                text = rewrite_image_paths(
                    chunk.get('text', '') or '',
                    output_dir,
                    args.public_prefix,
                ).strip()
                raw = page.get_text('text').strip()
                is_scanned = bool(page_profiles[page_number - 1]['is_scanned'])

                if is_scanned:
                    ocr_text, success, method = try_ocr(page, args.ocr_language, page_number)
                    if success:
                        text = f'{text}\n\n{ocr_text}'.strip()
                        ocr_done.append(page_number)
                        is_scanned = False
                        log_progress('ocr', page_number, total, f'OCR page {page_number} via {method}')
                    else:
                        log_progress('ocr', page_number, total, f'OCR impossible page {page_number}')

                if is_scanned:
                    scanned.append(page_number)

                raw_pages.append(text)

                if page_profiles[page_number - 1]['is_complex'] or page_requires_visual_reference(page, raw):
                    visual_queue.append(page_number)

    log_progress('visual', 0, max(len(visual_queue), 1),
                 f'Rendu de {len(visual_queue)} pages visuelles')

    # ── Phase 2: parallel visual rendering ───────────────────────────────
    if visual_queue:
        if parallel > 1:

            def render_one(pn: int) -> int:
                worker_doc = pymupdf.open(str(input_path))
                try:
                    render_page(
                        worker_doc[pn - 1],
                        output_dir / f'page-{pn:04d}.png',
                        args.visual_dpi,
                    )
                    return pn
                finally:
                    worker_doc.close()

            with ThreadPoolExecutor(max_workers=parallel) as ex:
                futs = {ex.submit(render_one, pn): pn for pn in visual_queue}
                for f in as_completed(futs):
                    try:
                        f.result()
                    except Exception as e:
                        warnings.append(f'Erreur rendu page {futs[f]}: {e}')
        else:
            for pn in visual_queue:
                try:
                    render_page(doc[pn - 1], output_dir / f'page-{pn:04d}.png', args.visual_dpi)
                except Exception as e:
                    warnings.append(f'Erreur rendu page {pn}: {e}')

    doc.close()

    # ── Phase 3: post-processing ─────────────────────────────────────────
    log_progress('clean', 0, len(raw_pages), 'Nettoyage et assemblage')
    boundaries = detect_repeated_boundaries(raw_pages)
    cleaned = [clean_markdown_page(p, boundaries) for p in raw_pages]
    cleaned = [p for p in cleaned if p]

    # Apply code detection and equation wrapping per-page (preserves code block structure)
    log_progress('enhance', 0, len(cleaned), 'Détection code et équations')
    enhanced = [post_process_markdown(p) for p in cleaned]

    markdown = '\n\n'.join(enhanced).strip()

    if scanned:
        warnings.append(
            'Certaines pages semblent scannées et nécessitent Tesseract OCR: '
            + ', '.join(str(p) for p in scanned) + '.'
        )

    images = sorted(
        p.name for p in output_dir.glob('*.png')
        if p.name != cover and not p.name.startswith('page-')
    )
    wc = len(re.findall(r'\b[\wÀ-ÿ\'-]+\b', markdown, flags=re.UNICODE))

    return {
        'markdown': markdown,
        'document_type': document_profile['document_type'],
        'extraction_strategy': document_profile['strategy'],
        'page_count': total,
        'word_count': wc,
        'image_count': len(images),
        'cover_file': cover,
        'visual_pages': visual_queue,
        'ocr_pages': ocr_done,
        'ocr_required_pages': scanned,
        'warnings': warnings,
    }

# ── Main dispatcher ──────────────────────────────────────────────────────────


def extract_document(args: argparse.Namespace) -> Dict[str, Any]:
    path = args.input.resolve()
    args.output_dir.resolve().mkdir(parents=True, exist_ok=True)

    if not path.is_file():
        raise ValueError(f'Fichier introuvable: {path}')
    if path.suffix.lower() not in SUPPORTED_EXTENSIONS:
        raise ValueError(
            f'Format non supporté: {path.suffix}. '
            f'Acceptés: {", ".join(sorted(SUPPORTED_EXTENSIONS))}'
        )

    ftype = detect_file_type(path)
    log_progress('detect', 0, 1, f'Type détecté: {ftype}')

    if ftype == 'office':
        return extract_office(args)
    return process_pdf(args)


def main(argv: Optional[List[str]] = None) -> int:
    args = parse_args(argv)
    try:
        result = extract_document(args)
    except Exception as e:
        log_progress('error', 0, 1, str(e))
        print(json.dumps({'error': str(e)}, ensure_ascii=False), file=sys.stderr)
        return 1
    print(json.dumps(result, ensure_ascii=False))
    return 0


if __name__ == '__main__':
    raise SystemExit(main())
